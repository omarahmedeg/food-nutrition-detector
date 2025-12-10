"""
Generate PowerPoint presentation for Food Nutrition Detector project
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_title_slide(prs, title, subtitle):
    """Create title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    
    # Style title
    title_frame = slide.shapes.title.text_frame
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    
    return slide

def create_content_slide(prs, title, content_items, layout_idx=1):
    """Create content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    slide.shapes.title.text = title
    
    # Style title
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(40)
    slide.shapes.title.text_frame.paragraphs[0].font.bold = True
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    
    # Add content
    body = slide.placeholders[1].text_frame
    body.clear()
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = body.paragraphs[0]
        else:
            p = body.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.level = 0
        p.space_before = Pt(6)
        
    return slide

def create_architecture_slide(prs):
    """Create architecture diagram slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "System Architecture"
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    
    # Architecture boxes
    # Frontend
    frontend = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(1.5), Inches(2), Inches(1)
    )
    frontend.fill.solid()
    frontend.fill.fore_color.rgb = RGBColor(135, 206, 250)
    frontend.line.color.rgb = RGBColor(0, 102, 204)
    frontend.line.width = Pt(2)
    text_frame = frontend.text_frame
    text_frame.text = "Next.js Frontend\n(Image Upload)"
    text_frame.paragraphs[0].font.size = Pt(14)
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Flask API
    api = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.5), Inches(1.5), Inches(2), Inches(1)
    )
    api.fill.solid()
    api.fill.fore_color.rgb = RGBColor(144, 238, 144)
    api.line.color.rgb = RGBColor(0, 128, 0)
    api.line.width = Pt(2)
    text_frame = api.text_frame
    text_frame.text = "Flask API\n(/api/predict)"
    text_frame.paragraphs[0].font.size = Pt(14)
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # PyTorch Model
    model = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.2), Inches(1.5), Inches(2.3), Inches(1)
    )
    model.fill.solid()
    model.fill.fore_color.rgb = RGBColor(255, 182, 193)
    model.line.color.rgb = RGBColor(220, 20, 60)
    model.line.width = Pt(2)
    text_frame = model.text_frame
    text_frame.text = "PyTorch Model\n(ResNet50 + Food-101)"
    text_frame.paragraphs[0].font.size = Pt(14)
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Nutrition DB
    nutrition = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.5), Inches(3.2), Inches(2), Inches(1)
    )
    nutrition.fill.solid()
    nutrition.fill.fore_color.rgb = RGBColor(255, 218, 185)
    nutrition.line.color.rgb = RGBColor(255, 140, 0)
    nutrition.line.width = Pt(2)
    text_frame = nutrition.text_frame
    text_frame.text = "Nutrition DB\n(101 class mappings)"
    text_frame.paragraphs[0].font.size = Pt(14)
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Response
    response = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(3.2), Inches(2), Inches(1)
    )
    response.fill.solid()
    response.fill.fore_color.rgb = RGBColor(221, 160, 221)
    response.line.color.rgb = RGBColor(147, 112, 219)
    response.line.width = Pt(2)
    text_frame = response.text_frame
    text_frame.text = "JSON Response\n(prediction + nutrition)"
    text_frame.paragraphs[0].font.size = Pt(14)
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Arrows
    # Frontend -> API
    arrow1 = slide.shapes.add_connector(1, Inches(2.8), Inches(2), Inches(3.5), Inches(2))
    arrow1.line.color.rgb = RGBColor(0, 0, 0)
    arrow1.line.width = Pt(2)
    
    # API -> Model
    arrow2 = slide.shapes.add_connector(1, Inches(5.5), Inches(2), Inches(6.2), Inches(2))
    arrow2.line.color.rgb = RGBColor(0, 0, 0)
    arrow2.line.width = Pt(2)
    
    # API -> Nutrition
    arrow3 = slide.shapes.add_connector(1, Inches(4.5), Inches(2.5), Inches(4.5), Inches(3.2))
    arrow3.line.color.rgb = RGBColor(0, 0, 0)
    arrow3.line.width = Pt(2)
    
    # API -> Response
    arrow4 = slide.shapes.add_connector(1, Inches(3.5), Inches(2), Inches(2.8), Inches(3.2))
    arrow4.line.color.rgb = RGBColor(0, 0, 0)
    arrow4.line.width = Pt(2)
    
    # Labels on arrows
    label1 = slide.shapes.add_textbox(Inches(2.9), Inches(1.7), Inches(0.5), Inches(0.3))
    label1.text_frame.text = "POST"
    label1.text_frame.paragraphs[0].font.size = Pt(10)
    label1.text_frame.paragraphs[0].font.italic = True
    
    return slide

def create_code_slide(prs, title, code_snippet):
    """Create slide with code snippet"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    
    # Code box
    code_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.2), Inches(9), Inches(4.5)
    )
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = RGBColor(43, 43, 43)
    code_box.line.color.rgb = RGBColor(100, 100, 100)
    
    text_frame = code_box.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = code_snippet
    p.font.name = 'Courier New'
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(220, 220, 220)
    
    return slide

def create_workflow_slide(prs):
    """Create workflow diagram slide showing the 5-step process"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "System Workflow"
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    
    # Define workflow steps
    steps = [
        {
            "number": "1",
            "title": "Upload Photo",
            "desc": "Frontend\n(User uploads image)",
            "color": RGBColor(135, 206, 250),
            "y": 1.5
        },
        {
            "number": "2",
            "title": "Send to Backend",
            "desc": "API Request\n(POST /api/predict)",
            "color": RGBColor(255, 182, 193),
            "y": 2.6
        },
        {
            "number": "3",
            "title": "Model Classifies",
            "desc": "PyTorch ResNet50\n(101 food classes)",
            "color": RGBColor(255, 218, 185),
            "y": 3.7
        },
        {
            "number": "4",
            "title": "Map to Nutrition",
            "desc": "Nutrition Database\n(CSV lookup)",
            "color": RGBColor(144, 238, 144),
            "y": 4.8
        },
        {
            "number": "5",
            "title": "Return to Frontend",
            "desc": "JSON Response\n(Results display)",
            "color": RGBColor(221, 160, 221),
            "y": 5.9
        }
    ]
    
    x_box = 1.5
    box_width = 7
    box_height = 0.9
    
    for step in steps:
        # Step box
        step_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x_box), Inches(step["y"]), Inches(box_width), Inches(box_height)
        )
        step_box.fill.solid()
        step_box.fill.fore_color.rgb = step["color"]
        step_box.line.color.rgb = RGBColor(0, 0, 0)
        step_box.line.width = Pt(2)
        
        # Step number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x_box - 0.5), Inches(step["y"] + 0.2), Inches(0.5), Inches(0.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(0, 102, 204)
        circle.line.color.rgb = RGBColor(0, 0, 0)
        circle.line.width = Pt(2)
        
        # Number text
        num_frame = circle.text_frame
        num_frame.text = step["number"]
        num_frame.paragraphs[0].font.size = Pt(20)
        num_frame.paragraphs[0].font.bold = True
        num_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        num_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        num_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Step text
        text_frame = step_box.text_frame
        text_frame.clear()
        
        # Title
        p1 = text_frame.paragraphs[0]
        p1.text = step["title"]
        p1.font.size = Pt(18)
        p1.font.bold = True
        p1.alignment = PP_ALIGN.CENTER
        
        # Description
        p2 = text_frame.add_paragraph()
        p2.text = step["desc"]
        p2.font.size = Pt(14)
        p2.alignment = PP_ALIGN.CENTER
        
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Arrow (except for last step)
        if step["number"] != "5":
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                Inches(4.5), Inches(step["y"] + 0.95), Inches(0.5), Inches(0.5)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(100, 100, 100)
            arrow.line.color.rgb = RGBColor(0, 0, 0)
    
    return slide

def create_visual_data_slide(prs):
    """Create slide with nutrition data visualization"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Sample Prediction Output"
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    
    # Prediction box
    pred_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.3), Inches(4), Inches(1.2)
    )
    pred_box.fill.solid()
    pred_box.fill.fore_color.rgb = RGBColor(240, 248, 255)
    pred_box.line.color.rgb = RGBColor(0, 102, 204)
    pred_box.line.width = Pt(2)
    
    text_frame = pred_box.text_frame
    text_frame.text = "Predicted: Pizza\nConfidence: 87%"
    for p in text_frame.paragraphs:
        p.font.size = Pt(20)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Nutrition info boxes
    nutrients = [
        ("Calories", "266 kcal", RGBColor(255, 99, 71)),
        ("Protein", "12g", RGBColor(60, 179, 113)),
        ("Carbs", "33g", RGBColor(255, 165, 0)),
        ("Fat", "10g", RGBColor(147, 112, 219)),
        ("Fiber", "2g", RGBColor(70, 130, 180)),
        ("Sugar", "4g", RGBColor(255, 105, 180))
    ]
    
    x_start = 0.5
    y_start = 2.8
    box_width = 1.4
    box_height = 0.8
    x_spacing = 1.5
    
    for i, (label, value, color) in enumerate(nutrients):
        row = i // 3
        col = i % 3
        
        x = x_start + (col * x_spacing)
        y = y_start + (row * 1)
        
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(box_width), Inches(box_height)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.color.rgb = RGBColor(255, 255, 255)
        box.line.width = Pt(2)
        
        text_frame = box.text_frame
        text_frame.text = f"{label}\n{value}"
        text_frame.paragraphs[0].font.size = Pt(16)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Add second paragraph for value
        if len(text_frame.paragraphs) > 1:
            text_frame.paragraphs[1].font.size = Pt(18)
            text_frame.paragraphs[1].font.bold = True
            text_frame.paragraphs[1].font.color.rgb = RGBColor(255, 255, 255)
            text_frame.paragraphs[1].alignment = PP_ALIGN.CENTER
    
    # Add JSON example on right side
    json_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5), Inches(1.3), Inches(4.5), Inches(4)
    )
    json_box.fill.solid()
    json_box.fill.fore_color.rgb = RGBColor(43, 43, 43)
    json_box.line.color.rgb = RGBColor(100, 100, 100)
    
    json_text = '''{
  "prediction": {
    "food_class": "pizza",
    "confidence": 0.87,
    "display_name": "Pizza"
  },
  "nutrition": {
    "calories": 266,
    "protein": 12,
    "carbs": 33,
    "fat": 10,
    "fiber": 2,
    "sugar": 4
  },
  "top_predictions": [...]
}'''
    
    text_frame = json_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = json_text
    p.font.name = 'Courier New'
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(144, 238, 144)
    
    return slide

def main():
    """Generate the full presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    create_title_slide(
        prs,
        "Food Nutrition Detector",
        "AI-Powered Food Classification with Nutritional Analysis\nDecember 8, 2025"
    )
    
    # Slide 2: Problem Statement
    create_content_slide(prs, "Problem Statement", [
        "Manual nutrition tracking is time-consuming and error-prone",
        "Users need quick, accurate calorie and nutrition estimates from photos",
        "Use cases: diet apps, health monitoring, restaurant menus, accessibility",
        "Solution: Automated food classification + nutrition lookup"
    ])
    
    # Slide 3: Project Goals
    create_content_slide(prs, "Project Goals", [
        "Use pre-trained Food-101 classifier (no training required)",
        "Provide top-K predictions with confidence scores",
        "Return detailed nutrition facts: calories, protein, carbs, fat, fiber, sugar",
        "Expose clean REST API for easy integration",
        "Simple, responsive web frontend for image upload"
    ])
    
    # Slide 4: Architecture
    create_architecture_slide(prs)
    
    # Slide 5: Workflow Diagram
    create_workflow_slide(prs)
    
    # Slide 6: Technology Stack
    create_content_slide(prs, "Technology Stack", [
        "Backend: Python 3.12 + Flask REST API",
        "ML Framework: PyTorch (ResNet50 architecture)",
        "Model: Food-101 pre-trained classifier (101 food classes)",
        "Frontend: Next.js + React",
        "Data: Custom nutrition database with per-class mappings",
        "Deployment: Local development server (production-ready architecture)"
    ])
    
    # Slide 7: Model Details
    create_content_slide(prs, "Model & Dataset", [
        "Pre-trained Food-101 checkpoint (food101_classifier.pth)",
        "Architecture: ResNet50 with modified final layer (101 classes)",
        "Dataset: Food-101 (101 common food categories)",
        "Inference-only: No training required",
        "Transform pipeline: Resize(224) → CenterCrop → Normalize",
        "Average inference time: <1 second on CPU"
    ])
    
    # Slide 8: Nutrition Database
    create_content_slide(prs, "Nutrition Database Design", [
        "Custom mapping: 101 food classes → nutrition facts",
        "Per-item data: calories, protein, carbs, fat, fiber, sugar",
        "Based on typical serving sizes (USDA/standardized portions)",
        "Fallback defaults for unknown items",
        "Extensible: easy to add vitamins, minerals, allergens",
        "File: nutrition_data_food101.py (Python dict with lookup function)"
    ])
    
    # Slide 9: API Specification
    create_content_slide(prs, "REST API Endpoints", [
        "POST /api/predict - Upload image, get prediction + nutrition",
        "  • Input: multipart/form-data (image file)",
        "  • Output: JSON with prediction, confidence, nutrition, top-K list",
        "GET /api/classes - List all supported food classes (101 items)",
        "GET /api/health - Health check endpoint",
        "Response time: typically <1s for classification + lookup",
        "Error handling: validation, fallbacks, logging"
    ])
    
    # Slide 10: Sample Output (Visual)
    create_visual_data_slide(prs)
    
    # Slide 11: PyTorch Model Code
    create_code_slide(prs, "PyTorch Model Loader (Excerpt)", '''class FastFoodClassifier:
    def __init__(self, model_path, num_classes=101):
        self.device = torch.device('cuda' if torch.cuda.is_available() else 'cpu')
        
        # Load checkpoint (handles multiple formats)
        checkpoint = torch.load(model_path, map_location=self.device)
        
        if isinstance(checkpoint, dict) and 'model_state_dict' in checkpoint:
            state_dict = checkpoint['model_state_dict']
            self.class_names = checkpoint.get('labels', FOOD101_CLASSES)
            num_classes = checkpoint.get('num_classes', 101)
        else:
            state_dict = checkpoint
            self.class_names = FOOD101_CLASSES
        
        # Create ResNet50 and load weights
        self.model = models.resnet50(pretrained=False)
        self.model.fc = nn.Linear(self.model.fc.in_features, num_classes)
        self.model.load_state_dict(state_dict)
        self.model.to(self.device)
        self.model.eval()''')
    
    # Slide 12: API Prediction Flow
    create_code_slide(prs, "API Prediction Endpoint (Simplified)", '''@app.route('/api/predict', methods=['POST'])
def predict():
    file = request.files.get('image')
    img = Image.open(file.stream).convert('RGB')
    
    # Get top-5 predictions
    predictions = model.predict(img, top_k=5)
    
    # Map each prediction to nutrition info
    top_pred = predictions[0]
    nutrition = get_nutrition_info(top_pred['class'])
    
    top_predictions_with_nutrition = [
        {
            'food_class': pred['class'],
            'confidence': pred['confidence'],
            'nutrition': get_nutrition_info(pred['class'])
        }
        for pred in predictions
    ]
    
    return jsonify({
        'prediction': top_pred,
        'nutrition': nutrition,
        'top_predictions': top_predictions_with_nutrition
    })''')
    
    # Slide 13: Frontend Features
    create_content_slide(prs, "Web Frontend (Next.js)", [
        "Image upload component with drag-and-drop",
        "Real-time preview of uploaded image",
        "Display top prediction with confidence percentage",
        "Nutrition card showing all macro/micro nutrients",
        "Top-K predictions list (alternative suggestions)",
        "Responsive design for mobile and desktop",
        "Simple, intuitive UX"
    ])
    
    # Slide 14: Key Implementation Challenges
    create_content_slide(prs, "Challenges & Solutions", [
        "Challenge: Checkpoint format variations",
        "  → Solution: Robust loader handles dict/raw state_dict",
        "Challenge: Missing class labels in checkpoint",
        "  → Solution: Fallback Food-101 class list built-in",
        "Challenge: HuggingFace rate-limiting during prototyping",
        "  → Solution: Switched to local model, no external dependencies",
        "Challenge: Nutrition accuracy per serving size",
        "  → Solution: Document typical portions, allow future user input"
    ])
    
    # Slide 15: Deployment & Setup
    create_content_slide(prs, "Running the Project", [
        "Prerequisites: Python 3.8+, Node.js 16+",
        "Backend setup:",
        "  pip install -r ml-model/requirements.txt",
        "  python api_server.py  # Runs on http://localhost:5000",
        "Frontend setup:",
        "  npm install && npm run dev",
        "Model file: Place food101_classifier.pth in ml-model/saved_models/",
        "Full setup time: ~5 minutes"
    ])
    
    # Slide 16: Testing & Validation
    create_content_slide(prs, "Testing Strategy", [
        "Unit tests: Model loader, nutrition DB lookups",
        "Integration tests: API endpoints with sample images",
        "Manual testing: Upload diverse food images, verify predictions",
        "Performance: Measure inference latency (<1s target)",
        "Error handling: Invalid files, missing model, edge cases",
        "Future: Add CI/CD pipeline for automated testing"
    ])
    
    # Slide 17: Results & Performance
    create_content_slide(prs, "Results & Observations", [
        "Inference speed: ~0.5-1.0 seconds per image (CPU)",
        "Classification accuracy: Depends on Food-101 checkpoint quality",
        "Nutrition lookup: 100% coverage for all 101 classes",
        "API response time: <1 second end-to-end",
        "User feedback: Nutrition display greatly improves perceived value",
        "Successful predictions logged: 10+ test requests with 200 OK"
    ])
    
    # Slide 18: Security & Privacy
    create_content_slide(prs, "Security & Privacy Considerations", [
        "Images not stored by default (privacy-first)",
        "HTTPS required for production deployment",
        "Input validation on all file uploads",
        "Rate limiting recommended for public APIs",
        "Consider on-device inference for sensitive use cases",
        "Data retention policy: configurable, opt-in only"
    ])
    
    # Slide 19: Limitations & Future Work
    create_content_slide(prs, "Limitations & Known Issues", [
        "Model limited to 101 Food-101 classes",
        "Nutrition values are approximations (typical servings)",
        "No portion size detection (assumes standard serving)",
        "Ambiguous multi-food images may confuse classifier",
        "No allergen or ingredient-level breakdown yet",
        "Recommendation: Allow manual corrections in UI"
    ])
    
    # Slide 20: Roadmap - Short Term
    create_content_slide(prs, "Roadmap: Next Steps (Short Term)", [
        "Add serving size input field (multiply nutrition values)",
        "Expand nutrition DB with vitamins, minerals, allergens",
        "Implement user feedback loop (correct wrong predictions)",
        "Add confidence threshold warnings (low-confidence predictions)",
        "Docker containerization for easy deployment",
        "Deploy to cloud (Azure/GCP/AWS) with CI/CD"
    ])
    
    # Slide 21: Roadmap - Long Term
    create_content_slide(prs, "Roadmap: Future Enhancements (Long Term)", [
        "Multi-food detection (object detection + segmentation)",
        "Mobile app with on-device inference (CoreML/TensorFlow Lite)",
        "Meal planning integration (track daily intake)",
        "Barcode scanning for packaged foods",
        "Recipe suggestions based on detected ingredients",
        "Fine-tune model on user-submitted corrections"
    ])
    
    # Slide 22: Code Repository
    create_content_slide(prs, "Project Structure", [
        "Repository: food-nutrition-detector/",
        "  api_server.py - Flask REST API (main entry point)",
        "  ml-model/",
        "    • pytorch_model.py - Model loader & predictor",
        "    • nutrition_data_food101.py - Nutrition mappings",
        "    • saved_models/ - Model checkpoints",
        "  Frontend files - Next.js/React components",
        "  generate_presentation.py - This presentation generator!",
        "Documentation: README.md with full setup instructions"
    ])
    
    # Slide 23: Demo Video / Live Demo
    create_content_slide(prs, "Live Demo", [
        "Upload an image of food (e.g., pizza, sushi, salad)",
        "API classifies the image in <1 second",
        "Frontend displays:",
        "  • Top prediction with confidence %",
        "  • Large calorie number",
        "  • Nutrition breakdown (protein, carbs, fat, fiber, sugar)",
        "  • Alternative top-3 predictions",
        "[Insert screenshot or live demo here]"
    ])
    
    # Slide 24: Lessons Learned
    create_content_slide(prs, "Lessons Learned", [
        "Robust checkpoint loading is critical for model portability",
        "Nutrition data quality directly impacts user trust",
        "Simple API design accelerates frontend integration",
        "Pre-trained models save weeks of training time",
        "User feedback loops improve accuracy over time",
        "Documentation & testing prevent deployment headaches"
    ])
    
    # Slide 25: Acknowledgments
    create_content_slide(prs, "Acknowledgments & Resources", [
        "Food-101 dataset creators (ETH Zurich)",
        "PyTorch & torchvision teams",
        "Flask & Next.js communities",
        "USDA nutrition database for reference values",
        "OpenAI for development assistance",
        "Community testers for feedback"
    ])
    
    # Slide 26: Thank You / Q&A
    create_title_slide(
        prs,
        "Thank You!",
        "Questions & Discussion\n\nRepository: food-nutrition-detector\nDemo: http://localhost:5000"
    )
    
    # Save presentation
    output_path = r"c:\Users\omara\Desktop\cv\food-nutrition-detector\Food_Nutrition_Detector_Presentation.pptx"
    prs.save(output_path)
    print(f"✓ Presentation created successfully!")
    print(f"✓ Saved to: {output_path}")
    print(f"✓ Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
